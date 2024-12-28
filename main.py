# Install necessary packages first
# pip install pytube opencv-python reportlab python-pptx

# from pytube import YouTube
from pytubefix import YouTube
from pytubefix.cli import on_progress
import cv2
import os
from reportlab.lib.pagesizes import landscape
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches
import shutil

# Function to download YouTube video using pytube
def download_video(url, output_path="downloads"):
    yt = YouTube(url, on_progress_callback = on_progress)    
    stream = yt.streams.get_highest_resolution()
    stream.download(output_path)
    return f"{output_path}/{yt.title}.mp4"

    

# Function to extract frames from the video using opencv
def extract_frames(video_path, output_folder="frames", fps=1):
    # Create folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Open video
    video = cv2.VideoCapture(video_path)
    frame_rate = video.get(cv2.CAP_PROP_FPS)  # Get the frames per second
    frame_interval = int(frame_rate / fps)  # Interval for frame extraction

    frame_count = 0
    saved_frame_count = 0

    while video.isOpened():
        ret, frame = video.read()
        if not ret:
            break
        frame_count += 1
        
        # Extract frames at the given interval
        if frame_count % frame_interval == 0:
            frame_filename = f"{output_folder}/frame_{saved_frame_count:04d}.png"
            cv2.imwrite(frame_filename, frame)  # Save frame as image
            saved_frame_count += 1

    video.release()
    print(f"Frames saved to {output_folder}")

# Function to convert images to PDF using reportlab
def images_to_pdf(image_folder, output_pdf="output.pdf"):
    c = canvas.Canvas("output.pdf", pagesize=landscape((1280, 720)))

    image_files = sorted([f for f in os.listdir(image_folder) if f.endswith('.png')])

    for img_file in image_files:
        img_path = os.path.join(image_folder, img_file)
        c.drawImage(img_path, 0, 0, width=1280, height=720)  # Use full resolution
        c.showPage()   # Add a new page for each image

    c.save()
    print(f"PDF saved as {output_pdf}")

# Function to convert images to PowerPoint using python-pptx
def images_to_ppt(image_folder, output_ppt="output.pptx"):
    prs = Presentation()
    image_files = sorted([f for f in os.listdir(image_folder) if f.endswith('.png')])

    for img_file in image_files:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 5 is a blank layout
        img_path = os.path.join(image_folder, img_file)
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.5), width=Inches(9))

    prs.save(output_ppt)
    print(f"PPT saved as {output_ppt}")

def clean_up():
    folders_to_delete = ["downloads", "frames"]

    for folder in folders_to_delete:
        folder_path = os.path.join(os.getcwd(), folder)  # Construct the full path
        try:
            shutil.rmtree(folder_path)  # Delete the folder and its contents
            print(f"Folder '{folder}' and all its contents have been deleted.")
        except FileNotFoundError:
            print(f"Folder '{folder}' not found.")
        except PermissionError:
            print(f"Permission denied to delete '{folder}'.")
        except Exception as e:
            print(f"An error occurred while deleting '{folder}': {e}")

# Complete workflow: Download video, extract frames, convert to PDF and PPT
def convert_video_to_pdf_and_ppt(video_url):
    # Step 1: Download video
    video_file = download_video(video_url)

    # Step 2: Extract frames from the video
    extract_frames(video_file)

    # Step 3: Convert extracted frames to PDF
    images_to_pdf("frames")

    # # Step 4: Convert extracted frames to PowerPoint
    # images_to_ppt("frames")

    clean_up()

# Example usage
if __name__ == "__main__":
    video_url = "https://www.youtube.com/watch?v=4xG2aJa6UyY"  # Replace with actual URL
    convert_video_to_pdf_and_ppt(video_url)